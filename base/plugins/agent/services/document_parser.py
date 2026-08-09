"""
文档解析服务 - 支持多种格式
PDF, DOCX, XLSX, PPTX, TXT, MD 等
"""
import logging
from pathlib import Path
from typing import Optional, Tuple
import io
from fastapi import UploadFile

logger = logging.getLogger(__name__)


class DocumentParser:
    """文档解析器"""
    
    @staticmethod
    async def parse_file(file: UploadFile) -> Tuple[str, str]:
        """解析上传的文件，返回（文本内容,文件类型）"""
        filename = file.filename or "unknown"
        content = await file.read()
        file_ext = Path(filename).suffix.lower()
        
        try:
            if file_ext in ['.pdf']:
                return DocumentParser._parse_pdf(content), 'pdf'
            elif file_ext in ['.docx']:
                return DocumentParser._parse_docx(content), 'docx'
            elif file_ext in ['.xlsx', '.xls']:
                return DocumentParser._parse_xlsx(content), 'xlsx'
            elif file_ext in ['.pptx']:
                return DocumentParser._parse_pptx(content), 'pptx'
            elif file_ext in ['.txt']:
                return DocumentParser._parse_txt(content), 'txt'
            elif file_ext in ['.md', '.markdown']:
                return DocumentParser._parse_txt(content), 'md'
            elif file_ext in ['.py', '.js', '.html', '.css', '.json', '.yaml', '.yml']:
                return DocumentParser._parse_txt(content), file_ext[1:]
            else:
                # 默认尝试以文本读取
                try:
                    return DocumentParser._parse_txt(content), file_ext[1:] if file_ext else 'txt'
                except:
                    return "", file_ext[1:] if file_ext else 'txt'
        except Exception as e:
            logger.error(f"解析文件 {filename} 失败: {e}")
            raise ValueError(f"不支持的文件格式或解析失败: {file_ext}")
    
    @staticmethod
    def _parse_pdf(content: bytes) -> str:
        """解析 PDF 文档"""
        try:
            import pdfplumber
            
            text_parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        text_parts.append(f"--- 第 {page_num + 1} 页 ---\n{page_text}")
            
            full_text = "\n\n".join(text_parts)
            
            if not full_text.strip():
                # 如果没有提取到文本，尝试用 PyPDF
                try:
                    from pypdf import PdfReader
                    reader = PdfReader(io.BytesIO(content))
                    text_parts = []
                    for page_num, page in enumerate(reader.pages):
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"--- 第 {page_num + 1} 页 ---\n{page_text}")
                    full_text = "\n\n".join(text_parts)
                except Exception as e:
                    logger.warning(f"PyPDF 也失败: {e}")
            
            return full_text
        except ImportError:
            raise ValueError("请安装 pdfplumber: pip install pdfplumber")
    
    @staticmethod
    def _parse_docx(content: bytes) -> str:
        """解析 Word 文档"""
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # 也处理表格
            for table_num, table in enumerate(doc.tables):
                table_text = f"--- 表格 {table_num + 1} ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        table_text += row_text + "\n"
                text_parts.append(table_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("请安装 python-docx: pip install python-docx")
    
    @staticmethod
    def _parse_xlsx(content: bytes) -> str:
        """解析 Excel 文档"""
        try:
            import pandas as pd
            
            excel_file = pd.ExcelFile(io.BytesIO(content))
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---\n")
                
                # 转换为文本
                csv_text = df.to_csv(sep=' | ', na_rep='', index=False, header=False)
                text_parts.append(csv_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("请安装 pandas 和 openpyxl: pip install pandas openpyxl")
    
    @staticmethod
    def _parse_pptx(content: bytes) -> str:
        """解析 PPT 文档"""
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"--- 第 {slide_num + 1} 页 ---\n"
                
                # 提取文本框
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                # 提取备注
                if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    slide_text += f"\n备注: {slide.notes_slide.notes_text_frame.text}\n"
                
                if slide_text.strip():
                    text_parts.append(slide_text)
            
            return "\n\n".join(text_parts)
        except ImportError:
            raise ValueError("请安装 python-pptx: pip install python-pptx")
    
    @staticmethod
    def _parse_txt(content: bytes) -> str:
        """解析文本文件"""
        try:
            # 尝试 UTF-8 解码
            return content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                # 尝试 GBK 解码
                return content.decode('gbk')
            except:
                # 最后尝试用 latin-1
                return content.decode('latin-1', errors='ignore')


# 支持的文件扩展名
SUPPORTED_EXTENSIONS = [
    '.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.ppt',
    '.txt', '.md', '.markdown', '.py', '.js', '.html', '.css',
    '.json', '.yaml', '.yml'
]


def is_supported_file(filename: str) -> bool:
    """检查是否支持该文件格式"""
    if not filename:
        return False
    ext = Path(filename).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS
